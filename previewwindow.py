import os
import tempfile
from PyQt5.QtGui import QTextDocument, QPixmap
from PyQt5.QtWidgets import (QPushButton, QVBoxLayout, QWidget, QLabel, QSplitter, QComboBox,
                             QAction, QMessageBox, QMenu, QTextEdit, QDialog,
                             QDialogButtonBox, QLineEdit, QProgressBar)
from PyQt5.QtCore import Qt, QTimer, QUrl
import pdfplumber
from docx import Document
from pptx import Presentation
import fitz  # 使用最新版 PyMuPDF 1.23.0 或更高版本
import shutil







class PreviewWindow(QWidget):
    """文件预览窗口 - 使用安全预览模式"""

    def __init__(self):
        super().__init__()
        self.setWindowTitle("文件预览")
        self.setGeometry(100, 100, 800, 600)

        self.layout = QVBoxLayout()

        # 状态标签
        self.status_label = QLabel("准备预览...")
        self.status_label.setAlignment(Qt.AlignCenter)
        self.layout.addWidget(self.status_label)

        # 进度条
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        self.layout.addWidget(self.progress_bar)

        # 安全预览模式按钮
        self.safe_mode_btn = QPushButton("启用安全预览模式")
        self.safe_mode_btn.setCheckable(True)
        self.safe_mode_btn.setChecked(True)
        self.safe_mode_btn.toggled.connect(self.toggle_safe_mode)
        self.layout.addWidget(self.safe_mode_btn)

        # 文本预览区域
        self.preview_text = QTextEdit()
        self.preview_text.setReadOnly(True)
        self.layout.addWidget(self.preview_text)

        self.setLayout(self.layout)
        self.safe_mode = True

    def toggle_safe_mode(self, enabled):
        """切换安全预览模式"""
        self.safe_mode = enabled
        mode_text = "启用" if enabled else "禁用"
        self.status_label.setText(f"安全预览模式已{mode_text}")

    def show_pdf_preview(self, filepath):
        """安全预览PDF文件"""
        if self.safe_mode:
            self._show_pdf_text_preview(filepath)
        else:
            self._show_pdf_image_preview(filepath)

    def _show_pdf_text_preview(self, filepath):
        """纯文本预览模式"""
        self.status_label.setText("正在提取PDF文本内容...")
        self.progress_bar.setVisible(True)
        self.progress_bar.setRange(0, 0)  # 不确定进度

        QTimer.singleShot(100, lambda: self._load_pdf_text_content(filepath))

    def _load_pdf_text_content(self, filepath):
        """加载PDF文本内容"""
        try:
            content = ""
            # 使用双重解析机制
            try:
                # 首先尝试使用 PyMuPDF
                with fitz.open(filepath) as doc:
                    total_pages = len(doc)
                    self.progress_bar.setRange(0, total_pages)

                    for i in range(min(5, total_pages)):  # 只预览前5页
                        page = doc.load_page(i)
                        text = page.get_text("text")
                        if text:
                            content += f"=== 第 {i + 1} 页 ===\n{text}\n\n"
                        else:
                            content += f"=== 第 {i + 1} 页 (无文本内容) ===\n\n"
                        self.progress_bar.setValue(i + 1)
            except Exception as fitz_error:
                # 失败时使用 pdfplumber
                with pdfplumber.open(filepath) as pdf:
                    total_pages = len(pdf.pages)
                    self.progress_bar.setRange(0, total_pages)

                    for i, page in enumerate(pdf.pages[:5]):  # 只预览前5页
                        try:
                            text = page.extract_text()
                            if text:
                                content += f"=== 第 {i + 1} 页 ===\n{text}\n\n"
                            else:
                                content += f"=== 第 {i + 1} 页 (无文本内容) ===\n\n"
                            self.progress_bar.setValue(i + 1)
                        except Exception as page_error:
                            content += f"=== 第 {i + 1} 页 (提取错误: {str(page_error)}) ===\n\n"

            if not content:
                content = "无法提取PDF文本内容"

            self.preview_text.setPlainText(content)
            self.status_label.setText(f"PDF预览: {os.path.basename(filepath)}")

        except Exception as e:
            self.preview_text.setPlainText(f"PDF加载失败: {str(e)}")
            self.status_label.setText("预览错误")

        finally:
            self.progress_bar.setVisible(False)

    def _show_pdf_image_preview(self, filepath):
        """图像预览模式"""
        self.status_label.setText("正在渲染PDF预览...")
        self.progress_bar.setVisible(True)
        self.progress_bar.setRange(0, 0)

        QTimer.singleShot(100, lambda: self._render_pdf_image_preview(filepath))

    def _render_pdf_image_preview(self, filepath):
        """渲染PDF图像预览"""
        try:
            # 创建临时目录
            temp_dir = tempfile.mkdtemp()

            # 渲染第一页为图像
            page_index = 0
            zoom = 1.0
            image_path = os.path.join(temp_dir, "preview.png")

            # 使用PyMuPDF渲染预览图
            with fitz.open(filepath) as doc:
                if len(doc) == 0:
                    raise Exception("空PDF文档")

                # 获取页面
                page = doc.load_page(page_index)

                # 设置渲染矩阵
                mat = fitz.Matrix(zoom, zoom)

                # 渲染页面为Pixmap
                pix = page.get_pixmap(matrix=mat, alpha=False)

                # 保存为PNG
                pix.save(image_path)

            # 显示图像预览
            pixmap = QPixmap(image_path)
            if not pixmap.isNull():
                self.preview_text.clear()
                self.preview_text.document().addResource(
                    QTextDocument.ImageResource,
                    QUrl("image://preview"),
                    pixmap
                )
                cursor = self.preview_text.textCursor()
                cursor.insertHtml(f'<center><img src="image://preview"></center>')
                self.status_label.setText(f"PDF图像预览: {os.path.basename(filepath)}")
            else:
                self.preview_text.setPlainText("无法渲染图像预览")
                self.status_label.setText("渲染失败")

        except Exception as e:
            self.preview_text.setPlainText(f"图像预览失败: {str(e)}\n已自动切换到文本预览模式")
            self.safe_mode_btn.setChecked(True)  # 自动切换到安全模式
            self._show_pdf_text_preview(filepath)  # 回退到文本预览

        finally:
            # 清理临时文件
            try:
                shutil.rmtree(temp_dir)
            except:
                pass

            self.progress_bar.setVisible(False)

    # Word文档预览
    def show_docx_preview(self, filepath):
        """预览Word文档"""
        self.status_label.setText("正在加载Word文档...")
        self.progress_bar.setVisible(True)
        self.progress_bar.setRange(0, 0)

        QTimer.singleShot(100, lambda: self._load_docx_content(filepath))

    def _load_docx_content(self, filepath):
        """加载Word内容"""
        try:
            doc = Document(filepath)
            content = ""
            paragraphs = doc.paragraphs[:100]  # 限制前100段
            self.progress_bar.setRange(0, len(paragraphs))

            for i, para in enumerate(paragraphs):
                content += para.text + "\n"
                if i % 10 == 0:  # 每10段更新一次进度
                    self.progress_bar.setValue(i)

            if not content:
                content = "无法提取Word文档内容"

            self.preview_text.setPlainText(content)
            self.status_label.setText(f"Word预览: {os.path.basename(filepath)}")

        except Exception as e:
            self.preview_text.setPlainText(f"Word文件加载失败: {str(e)}")
            self.status_label.setText("预览错误")

        finally:
            self.progress_bar.setVisible(False)

    # PPT文档预览
    def show_pptx_preview(self, filepath):
        """预览PPT文档"""
        self.status_label.setText("正在加载PPT文档...")
        self.progress_bar.setVisible(True)
        self.progress_bar.setRange(0, 0)

        QTimer.singleShot(100, lambda: self._load_pptx_content(filepath))

    def _load_pptx_content(self, filepath):
        """加载PPT内容"""
        try:
            prs = Presentation(filepath)
            content = ""
            slides = prs.slides[:10]  # 限制前10页
            self.progress_bar.setRange(0, len(slides))

            for i, slide in enumerate(slides):
                slide_title = f"幻灯片 {i + 1}"
                if slide.shapes.title:
                    slide_title = slide.shapes.title.text

                content += f"=== {slide_title} ===\n"

                # 提取内容
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape != slide.shapes.title:
                        slide_content.append(shape.text)

                content += "\n".join(slide_content) + "\n\n"
                self.progress_bar.setValue(i + 1)

            if not content:
                content = "无法提取PPT内容"

            self.preview_text.setPlainText(content)
            self.status_label.setText(f"PPT预览: {os.path.basename(filepath)}")

        except Exception as e:
            self.preview_text.setPlainText(f"PPT文件加载失败: {str(e)}")
            self.status_label.setText("预览错误")

        finally:
            self.progress_bar.setVisible(False)