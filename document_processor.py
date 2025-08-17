import re
import fitz
import pdfplumber
from docx import Document
from pdfminer.layout import LAParams
from pptx import Presentation


class DocumentProcessor:
    """文档处理核心模块 - 完全移除 PyMuPDF 依赖"""

    @staticmethod
    def extract_pdf_content(filepath):
        """
        兼容最新 pdfplumber 版本的 PDF 解析方法
        """
        content = {"metadata": {}, "sections": []}

        # 优先使用 PyMuPDF (fitz) 解析器
        result = DocumentProcessor.extract_with_pymupdf(filepath)
        if result and result.get('sections') and result['sections'][0].get('content'):
            return result

        # 如果 PyMuPDF 失败，尝试 pdfplumber
        print("PyMuPDF 解析不完整，尝试 pdfplumber")
        return DocumentProcessor.extract_with_pdfplumber(filepath)

    @staticmethod
    def extract_with_pymupdf(filepath):
        """使用 PyMuPDF 作为主要解析器"""
        try:
            content = {"metadata": {}, "sections": []}
            current_section = None

            doc = fitz.open(filepath)
            content['metadata']['pages'] = len(doc)
            content['metadata']['author'] = doc.metadata.get('author', 'Unknown')

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                try:
                    # PyMuPDF 通常能更好地处理颜色空间问题
                    text = page.get_text("text")

                    if text:
                        # 清理页码和特殊字符
                        cleaned_text = DocumentProcessor.clean_content(text)

                        # 处理内容结构
                        lines = cleaned_text.split('\n')
                        skip_section = False
                        for line in lines:
                            stripped_line = line.strip()
                            if stripped_line:
                                # 检查是否为非概念性内容的标题
                                if DocumentProcessor.is_heading(stripped_line) and not DocumentProcessor.is_concept_content(stripped_line):
                                    skip_section = True
                                    current_section = None
                                    continue
                                
                                # 如果需要跳过当前章节，继续跳过
                                if skip_section:
                                    if DocumentProcessor.is_heading(stripped_line):
                                        # 新的标题出现，检查是否为概念性内容
                                        if DocumentProcessor.is_concept_content(stripped_line):
                                            skip_section = False
                                        else:
                                            continue
                                    else:
                                        continue
                                
                                # 如果还没有创建节，或者遇到可能的标题
                                if current_section is None or DocumentProcessor.is_heading(stripped_line):
                                    if current_section and current_section['content']:
                                        content['sections'].append(current_section)
                                    # 检查新章节是否为概念性内容
                                    if DocumentProcessor.is_concept_content(stripped_line):
                                        current_section = {
                                            'title': stripped_line,
                                            'content': []
                                        }
                                    else:
                                        current_section = None
                                else:
                                    if current_section is not None and DocumentProcessor.is_concept_content(stripped_line):
                                        current_section['content'].append(stripped_line)
                except Exception as e:
                    print(f"第 {page_num + 1} 页 PyMuPDF 解析错误: {str(e)}")
                    if current_section is None:
                        current_section = {"title": "解析错误", "content": []}
                    current_section['content'].append(f"第 {page_num + 1} 页解析失败")

            # 添加最后一节
            if current_section and current_section['content']:
                content['sections'].append(current_section)

            return content

        except Exception as e:
            print(f"PyMuPDF 完全失败: {str(e)}")
            return None

    @staticmethod
    def extract_with_pdfplumber(filepath):
        """使用 pdfplumber 解析，兼容最新版本"""
        content = {"metadata": {}, "sections": []}
        current_section = None

        try:
            # 使用 pdfminer 的 LAParams
            laparams = LAParams()

            with pdfplumber.open(filepath, laparams=laparams) as pdf:
                content['metadata']['pages'] = len(pdf.pages)
                content['metadata']['author'] = pdf.metadata.get('Author', 'Unknown')

                for i, page in enumerate(pdf.pages):
                    try:
                        # 尝试标准提取
                        text = page.extract_text()
                    except Exception as e:
                        # 处理颜色空间错误
                        if "invalid color" in str(e).lower() or "gray" in str(e).lower():
                            print(f"第 {i + 1} 页遇到颜色空间错误，尝试简化提取")
                            try:
                                # 尝试简化提取方法
                                text = page.extract_text_simple()
                            except:
                                text = ""
                                print(f"第 {i + 1} 页简化提取失败")
                        else:
                            text = ""
                            print(f"第 {i + 1} 页解析错误: {str(e)}")

                    if text:
                        # 清理页码
                        cleaned_text = DocumentProcessor.clean_content(text)

                        # 处理内容结构
                        lines = cleaned_text.split('\n')
                        skip_section = False
                        for line in lines:
                            stripped_line = line.strip()
                            if stripped_line:
                                # 检查是否为非概念性内容的标题
                                if DocumentProcessor.is_heading(stripped_line) and not DocumentProcessor.is_concept_content(stripped_line):
                                    skip_section = True
                                    current_section = None
                                    continue
                                
                                # 如果需要跳过当前章节，继续跳过
                                if skip_section:
                                    if DocumentProcessor.is_heading(stripped_line):
                                        # 新的标题出现，检查是否为概念性内容
                                        if DocumentProcessor.is_concept_content(stripped_line):
                                            skip_section = False
                                        else:
                                            continue
                                    else:
                                        continue
                                
                                if current_section is None or DocumentProcessor.is_heading(stripped_line):
                                    if current_section and current_section['content']:
                                        content['sections'].append(current_section)
                                    # 检查新章节是否为概念性内容
                                    if DocumentProcessor.is_concept_content(stripped_line):
                                        current_section = {
                                            'title': stripped_line,
                                            'content': []
                                        }
                                    else:
                                        current_section = None
                                else:
                                    if current_section is not None and DocumentProcessor.is_concept_content(stripped_line):
                                        current_section['content'].append(stripped_line)

            # 添加最后一节
            if current_section and current_section['content']:
                content['sections'].append(current_section)

            return content

        except Exception as e:
            print(f"pdfplumber 解析失败: {str(e)}")
            # 尝试使用 PyMuPDF 作为后备
            return DocumentProcessor.extract_with_pymupdf(filepath)

    @staticmethod
    def clean_content(text):
        """清理内容：去除页码和其他不需要的元素"""
        if not text:
            return ""

        # 1. 去除单独一行的数字（页码）
        cleaned = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)

        # 2. 去除类似 "P123" 的错误信息残留
        cleaned = re.sub(r'\bP\d+\b', '', cleaned)

        # 3. 去除常见的页码模式 (Page X, P.X, etc.)
        cleaned = re.sub(r'(?i)\b(?:page|p|pg|pag|pagina)\s*[\.:]*\s*\d+\b', '', cleaned)

        # 4. 去除多余空行
        cleaned = re.sub(r'\n\s*\n', '\n', cleaned)

        return cleaned.strip()

    @staticmethod
    def is_heading(text):
        """判断是否为标题的简单实现"""
        # 可以根据您的文档特点扩展此方法
        heading_keywords = ["目录", "章节", "节", "第", "摘要", "引言", "结论", "参考", "概述", "总结", "要点", "重点"]
        # 标题通常较短且包含关键词，或者以数字开头的章节标题
        if (len(text) < 50 and any(keyword in text for keyword in heading_keywords)) or \
           (len(text) < 30 and re.match(r"^[0-9]+[、. ]", text)):
            return True
        return False

    @staticmethod
    def is_concept_content(text):
        """判断是否为概念性内容"""
        # 排除例子、作业等非概念性内容的关键词
        non_concept_keywords = ["例", "例子", "示例", "例题", "习题", "练习", "作业", "课后练习", "思考题", "实验", "实训"]
        # 如果文本以非概念性关键词开头，则不是概念内容
        if any(text.strip().startswith(keyword) for keyword in non_concept_keywords):
            return False
        # 如果文本包含非概念性关键词且不是标题，则不是概念内容
        if any(keyword in text for keyword in non_concept_keywords) and not DocumentProcessor.is_heading(text):
            return False
        # 特殊处理：如果文本是数字加点开头的列表项，可能是练习题（除非是重点内容）
        if re.match(r"^\d+\.\s", text.strip()) and "重点" not in text and "要点" not in text:
            return False
        return True


    # 其他extract方法保持不变
    @staticmethod
    def extract_docx_content(filepath):
        """提取Word文档内容"""
        doc = Document(filepath)
        content = {"metadata": {}, "sections": []}
        current_section = None

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:  # 只处理非空段落
                if para.style.name.startswith('Heading') or DocumentProcessor.is_heading(text):
                    # 新章节开始
                    # 检查是否为非概念性内容的标题
                    if not DocumentProcessor.is_concept_content(text):
                        # 跳过这个章节的所有内容
                        current_section = None
                        continue
                    
                    if current_section and current_section['content']:
                        content['sections'].append(current_section)
                    current_section = {
                        'title': text,
                        'content': [text]
                    }
                else:
                    # 内容段落 - 只添加概念性内容
                    if current_section is not None and DocumentProcessor.is_concept_content(text):
                        current_section['content'].append(text)

        if current_section and current_section['content']:
            content['sections'].append(current_section)
        return content

    @staticmethod
    def extract_pptx_content(filepath):
        """提取PPT内容"""
        prs = Presentation(filepath)
        content = {"metadata": {}, "sections": []}
        current_section = None

        for i, slide in enumerate(prs.slides):
            slide_title = f"幻灯片 {i + 1}"
            slide_content = []

            # 提取幻灯片标题
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = slide.shapes.title.text.strip()

            # 提取内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape != slide.shapes.title and shape.text.strip():
                    # 尝试识别标题和内容结构
                    text = shape.text.strip()
                    if DocumentProcessor.is_heading(text):
                        # 检查是否为非概念性内容的标题
                        if not DocumentProcessor.is_concept_content(text):
                            # 跳过这个章节的所有内容
                            current_section = None
                            continue
                        
                        # 如果是标题，创建新章节
                        if current_section and current_section['content']:
                            content['sections'].append(current_section)
                        current_section = {
                            'title': text,
                            'content': []
                        }
                    else:
                        # 如果是内容，添加到当前章节 - 只添加概念性内容
                        if DocumentProcessor.is_concept_content(text) and current_section is not None:
                            current_section['content'].append(text)

        # 添加最后一节
        if current_section and current_section['content']:
            content['sections'].append(current_section)
            
        # 如果没有识别出结构，使用默认方式
        if not content['sections']:
            for i, slide in enumerate(prs.slides):
                slide_title = f"幻灯片 {i + 1}"
                slide_content = []

                # 提取幻灯片标题
                if slide.shapes.title and slide.shapes.title.text.strip():
                    slide_title = slide.shapes.title.text.strip()

                # 提取内容 - 只添加概念性内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape != slide.shapes.title and shape.text.strip():
                        text = shape.text.strip()
                        # 只添加概念性内容
                        if DocumentProcessor.is_concept_content(text):
                            slide_content.append(text)

                # 只有当slide_content不为空时才添加
                if slide_content:
                    content['sections'].append({
                        'title': slide_title,
                        'content': slide_content
                    })
                
        return content

    # generate_knowledge_tree 保持不变
    @staticmethod
    def generate_knowledge_tree(documents):
        knowledge_tree = []
        for doc in documents:
            for section in doc['sections']:
                node = {
                    'title': section['title'],
                    'content': '\n'.join(section['content']),
                    'children': []
                }
                knowledge_tree.append(node)
        return knowledge_tree

    # export_to_word 保持不变
    @staticmethod
    def export_to_word(knowledge_tree, output_path):
        doc = Document()
        
        # 统一设置默认字体和字号
        from docx.shared import Pt
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
        
        # 设置正文样式
        style = doc.styles['Normal']
        font = style.font
        font.name = '宋体'
        font.size = Pt(12)  # 设置默认字号为12磅
        
        # 确保中文字体正确应用到正文
        style.element.xpath('.//w:rPr')[0].insert(0, OxmlElement('w:rFonts'))
        style.element.xpath('.//w:rFonts')[0].set(qn('w:eastAsia'), '宋体')
        
        # 设置标题样式
        heading_style = doc.styles['Heading 1']
        heading_font = heading_style.font
        heading_font.name = '宋体'
        heading_font.size = Pt(16)  # 设置标题字号为16磅
        
        # 确保中文字体正确应用到标题
        heading_style.element.xpath('.//w:rPr')[0].insert(0, OxmlElement('w:rFonts'))
        heading_style.element.xpath('.//w:rFonts')[0].set(qn('w:eastAsia'), '宋体')
        
        for node in knowledge_tree:
            # 清理标题中的特殊字符
            clean_title = DocumentProcessor._clean_text_for_xml(node['title'])
            
            # 添加标题，统一格式
            heading = doc.add_heading(clean_title, level=1)
            heading.style = 'Heading 1'
            # 确保标题使用宋体
            for run in heading.runs:
                run.font.name = '宋体'
                run._element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
            
            # 统一内容段落格式
            if isinstance(node['content'], list):
                for item in node['content']:
                    # 清理内容中的特殊字符
                    clean_item = DocumentProcessor._clean_text_for_xml(item)
                    if clean_item.strip():  # 只添加非空内容
                        paragraph = doc.add_paragraph(clean_item)
                        paragraph.style = 'Normal'
                        # 统一字体和字号
                        for run in paragraph.runs:
                            run.font.name = '宋体'
                            run.font.size = Pt(12)
                            run._element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')
            else:
                # 清理内容中的特殊字符
                clean_content = DocumentProcessor._clean_text_for_xml(node['content'])
                if clean_content.strip():  # 只添加非空内容
                    paragraph = doc.add_paragraph(clean_content)
                    paragraph.style = 'Normal'
                    # 统一字体和字号
                    for run in paragraph.runs:
                        run.font.name = '宋体'
                        run.font.size = Pt(12)
                        run._element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')

        doc.save(output_path)
        return output_path

    @staticmethod
    def _clean_text_for_xml(text):
        """清理文本中的特殊字符，确保XML兼容性"""
        if not isinstance(text, str):
            text = str(text)
        
        # 移除控制字符（除了常见的换行符和制表符）
        cleaned = ''.join(char for char in text if ord(char) >= 32 or char in '\n\t')
        
        # 替换特殊XML字符
        cleaned = cleaned.replace('&', '&')
        cleaned = cleaned.replace('<', '<')
        cleaned = cleaned.replace('>', '>')
        
        return cleaned

    # export_to_pdf 保持不变（使用reportlab）
    @staticmethod
    def export_to_pdf(word_path, pdf_path):
        try:
            import comtypes.client
            import os

            # 确保路径为绝对路径
            word_path = os.path.abspath(word_path)
            pdf_path = os.path.abspath(pdf_path)

            # 创建Word应用实例（后台运行，不显示窗口）
            word = comtypes.client.CreateObject('Word.Application')
            word.Visible = False

            # 打开Word文档
            doc = word.Documents.Open(word_path)

            # 保存为PDF（FileFormat=17是Word中PDF的固定格式代码）
            doc.SaveAs2(pdf_path, FileFormat=17)

            # 关闭文档和Word进程
            doc.Close()
            word.Quit()

            return pdf_path if os.path.exists(pdf_path) else None
        except Exception as e:
            print(f"PDF导出失败: {str(e)}")
            return None
