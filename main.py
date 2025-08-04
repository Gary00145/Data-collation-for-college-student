import os
import sys
import tempfile
import subprocess

from PyQt5.QtGui import QTextDocument, QPixmap
from PyQt5.QtWidgets import (QApplication, QMainWindow, QFileDialog, QTreeWidget,
                             QTreeWidgetItem, QListWidget, QPushButton, QHBoxLayout,
                             QVBoxLayout, QWidget, QLabel, QSplitter, QComboBox,
                             QAction, QMessageBox, QMenu, QTextEdit, QDialog,
                             QDialogButtonBox, QLineEdit, QProgressBar)
from PyQt5.QtCore import Qt, QTimer, QUrl
import pdfplumber
from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
import fitz  # 使用最新版 PyMuPDF 1.23.0 或更高版本
import shutil
import re
from pdfminer.layout import LAParams
from docx.shared import Pt, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT, WD_LINE_SPACING
from docx2pdf import convert


class DocumentProcessor:
    """文档处理核心模块 - 完全移除 PyMuPDF 依赖"""

    @staticmethod
    def extract_pdf_content(filepath):
        """
        兼容最新 pdfplumber 版本的 PDF 解析方法
        """
        content = {"metadata": {}, "sections": []}

        # 优先使用 PyMuPDF (fitz) 解析器
        result = DocumentProcessor.extract_with_pymupdf(filepath)
        if result and result.get('sections') and result['sections'][0].get('content'):
            return result

        # 如果 PyMuPDF 失败，尝试 pdfplumber
        print("PyMuPDF 解析不完整，尝试 pdfplumber")
        return DocumentProcessor.extract_with_pdfplumber(filepath)

    @staticmethod
    def extract_with_pymupdf(filepath):
        """使用 PyMuPDF 作为主要解析器"""
        try:
            content = {"metadata": {}, "sections": []}
            current_section = None

            doc = fitz.open(filepath)
            content['metadata']['pages'] = len(doc)
            content['metadata']['author'] = doc.metadata.get('author', 'Unknown')

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                try:
                    # PyMuPDF 通常能更好地处理颜色空间问题
                    text = page.get_text("text")

                    if text:
                        # 清理页码和特殊字符
                        cleaned_text = DocumentProcessor.clean_content(text)

                        # 处理内容结构
                        lines = cleaned_text.split('\n')
                        for line in lines:
                            stripped_line = line.strip()
                            if stripped_line:
                                # 如果还没有创建节，或者遇到可能的标题
                                if current_section is None or DocumentProcessor.is_heading(stripped_line):
                                    if current_section and current_section['content']:
                                        content['sections'].append(current_section)
                                    current_section = {
                                        'title': stripped_line,
                                        'content': []
                                    }
                                else:
                                    current_section['content'].append(stripped_line)
                except Exception as e:
                    print(f"第 {page_num + 1} 页 PyMuPDF 解析错误: {str(e)}")
                    if current_section is None:
                        current_section = {"title": "解析错误", "content": []}
                    current_section['content'].append(f"第 {page_num + 1} 页解析失败")

            # 添加最后一节
            if current_section and current_section['content']:
                content['sections'].append(current_section)

            return content

        except Exception as e:
            print(f"PyMuPDF 完全失败: {str(e)}")
            return None

    @staticmethod
    def extract_with_pdfplumber(filepath):
        """使用 pdfplumber 解析，兼容最新版本"""
        content = {"metadata": {}, "sections": []}
        current_section = None

        try:
            # 使用 pdfminer 的 LAParams
            laparams = LAParams()

            with pdfplumber.open(filepath, laparams=laparams) as pdf:
                content['metadata']['pages'] = len(pdf.pages)
                content['metadata']['author'] = pdf.metadata.get('Author', 'Unknown')

                for i, page in enumerate(pdf.pages):
                    try:
                        # 尝试标准提取
                        text = page.extract_text()
                    except Exception as e:
                        # 处理颜色空间错误
                        if "invalid color" in str(e).lower() or "gray" in str(e).lower():
                            print(f"第 {i + 1} 页遇到颜色空间错误，尝试简化提取")
                            try:
                                # 尝试简化提取方法
                                text = page.extract_text_simple()
                            except:
                                text = ""
                                print(f"第 {i + 1} 页简化提取失败")
                        else:
                            text = ""
                            print(f"第 {i + 1} 页解析错误: {str(e)}")

                    if text:
                        # 清理页码
                        cleaned_text = DocumentProcessor.clean_content(text)

                        # 处理内容结构
                        lines = cleaned_text.split('\n')
                        for line in lines:
                            stripped_line = line.strip()
                            if stripped_line:
                                if current_section is None or DocumentProcessor.is_heading(stripped_line):
                                    if current_section and current_section['content']:
                                        content['sections'].append(current_section)
                                    current_section = {
                                        'title': stripped_line,
                                        'content': []
                                    }
                                else:
                                    if current_section is not None:
                                        current_section['content'].append(stripped_line)

            # 添加最后一节
            if current_section and current_section['content']:
                content['sections'].append(current_section)

            return content

        except Exception as e:
            print(f"pdfplumber 解析失败: {str(e)}")
            # 尝试使用 PyMuPDF 作为后备
            return DocumentProcessor.extract_with_pymupdf(filepath)

    @staticmethod
    def clean_content(text):
        """清理内容：去除页码和其他不需要的元素"""
        if not text:
            return ""

        # 1. 去除单独一行的数字（页码）
        cleaned = re.sub(r'^\s*\d+\s*$', '', text, flags=re.MULTILINE)

        # 2. 去除类似 "P123" 的错误信息残留
        cleaned = re.sub(r'\bP\d+\b', '', cleaned)

        # 3. 去除常见的页码模式 (Page X, P.X, etc.)
        cleaned = re.sub(r'(?i)\b(?:page|p|pg|pag|pagina)\s*[\.:]*\s*\d+\b', '', cleaned)

        # 4. 去除多余空行
        cleaned = re.sub(r'\n\s*\n', '\n', cleaned)

        return cleaned.strip()

    @staticmethod
    def is_heading(text):
        """判断是否为标题的简单实现"""
        # 可以根据您的文档特点扩展此方法
        heading_keywords = ["目录", "章节", "节", "第", "摘要", "引言", "结论", "参考"]
        if len(text) < 50 and any(keyword in text for keyword in heading_keywords):
            return True
        return False


    # 其他extract方法保持不变
    @staticmethod
    def extract_docx_content(filepath):
        """提取Word文档内容"""
        doc = Document(filepath)
        content = {"metadata": {}, "sections": []}
        current_section = {"title": "引言", "content": []}

        for para in doc.paragraphs:
            if para.style.name.startswith('Heading'):
                # 新章节开始
                if current_section['content']:
                    content['sections'].append(current_section)
                current_section = {
                    'title': para.text.strip(),
                    'content': [para.text]
                }
            else:
                current_section['content'].append(para.text)

        if current_section['content']:
            content['sections'].append(current_section)
        return content

    @staticmethod
    def extract_pptx_content(filepath):
        """提取PPT内容"""
        prs = Presentation(filepath)
        content = {"metadata": {}, "sections": []}

        for i, slide in enumerate(prs.slides):
            slide_title = f"幻灯片 {i + 1}"
            slide_content = []

            # 提取幻灯片标题
            if slide.shapes.title:
                slide_title = slide.shapes.title.text

            # 提取内容
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape != slide.shapes.title:
                    slide_content.append(shape.text)

            content['sections'].append({
                'title': slide_title,
                'content': '\n'.join(slide_content)
            })
        return content

    # generate_knowledge_tree 保持不变
    @staticmethod
    def generate_knowledge_tree(documents):
        knowledge_tree = []
        for doc in documents:
            for section in doc['sections']:
                node = {
                    'title': section['title'],
                    'content': '\n'.join(section['content']),
                    'children': []
                }
                knowledge_tree.append(node)
        return knowledge_tree

    # export_to_word 保持不变
    @staticmethod
    def export_to_word(knowledge_tree, output_path):
        doc = Document()

        for node in knowledge_tree:
            doc.add_heading(node['title'], level=1)
            if isinstance(node['content'], list):
                for item in node['content']:
                    doc.add_paragraph(item)
            else:
                doc.add_paragraph(node['content'])

        doc.save(output_path)
        return output_path

    # export_to_pdf 保持不变（使用reportlab）
    @staticmethod
    def export_to_pdf(word_path, pdf_path):
        try:
            import comtypes.client
            import os

            # 确保路径为绝对路径
            word_path = os.path.abspath(word_path)
            pdf_path = os.path.abspath(pdf_path)

            # 创建Word应用实例（后台运行，不显示窗口）
            word = comtypes.client.CreateObject('Word.Application')
            word.Visible = False

            # 打开Word文档
            doc = word.Documents.Open(word_path)

            # 保存为PDF（FileFormat=17是Word中PDF的固定格式代码）
            doc.SaveAs2(pdf_path, FileFormat=17)

            # 关闭文档和Word进程
            doc.Close()
            word.Quit()

            return pdf_path if os.path.exists(pdf_path) else None
        except Exception as e:
            print(f"PDF导出失败: {str(e)}")
            return None

class PreviewWindow(QWidget):
    """文件预览窗口 - 使用安全预览模式"""

    def __init__(self):
        super().__init__()
        self.setWindowTitle("文件预览")
        self.setGeometry(100, 100, 800, 600)

        self.layout = QVBoxLayout()

        # 状态标签
        self.status_label = QLabel("准备预览...")
        self.status_label.setAlignment(Qt.AlignCenter)
        self.layout.addWidget(self.status_label)

        # 进度条
        self.progress_bar = QProgressBar()
        self.progress_bar.setVisible(False)
        self.layout.addWidget(self.progress_bar)

        # 安全预览模式按钮
        self.safe_mode_btn = QPushButton("启用安全预览模式")
        self.safe_mode_btn.setCheckable(True)
        self.safe_mode_btn.setChecked(True)
        self.safe_mode_btn.toggled.connect(self.toggle_safe_mode)
        self.layout.addWidget(self.safe_mode_btn)

        # 文本预览区域
        self.preview_text = QTextEdit()
        self.preview_text.setReadOnly(True)
        self.layout.addWidget(self.preview_text)

        self.setLayout(self.layout)
        self.safe_mode = True

    def toggle_safe_mode(self, enabled):
        """切换安全预览模式"""
        self.safe_mode = enabled
        mode_text = "启用" if enabled else "禁用"
        self.status_label.setText(f"安全预览模式已{mode_text}")

    def show_pdf_preview(self, filepath):
        """安全预览PDF文件"""
        if self.safe_mode:
            self._show_pdf_text_preview(filepath)
        else:
            self._show_pdf_image_preview(filepath)

    def _show_pdf_text_preview(self, filepath):
        """纯文本预览模式"""
        self.status_label.setText("正在提取PDF文本内容...")
        self.progress_bar.setVisible(True)
        self.progress_bar.setRange(0, 0)  # 不确定进度

        QTimer.singleShot(100, lambda: self._load_pdf_text_content(filepath))

    def _load_pdf_text_content(self, filepath):
        """加载PDF文本内容"""
        try:
            content = ""
            # 使用双重解析机制
            try:
                # 首先尝试使用 PyMuPDF
                with fitz.open(filepath) as doc:
                    total_pages = len(doc)
                    self.progress_bar.setRange(0, total_pages)

                    for i in range(min(5, total_pages)):  # 只预览前5页
                        page = doc.load_page(i)
                        text = page.get_text("text")
                        if text:
                            content += f"=== 第 {i + 1} 页 ===\n{text}\n\n"
                        else:
                            content += f"=== 第 {i + 1} 页 (无文本内容) ===\n\n"
                        self.progress_bar.setValue(i + 1)
            except Exception as fitz_error:
                # 失败时使用 pdfplumber
                with pdfplumber.open(filepath) as pdf:
                    total_pages = len(pdf.pages)
                    self.progress_bar.setRange(0, total_pages)

                    for i, page in enumerate(pdf.pages[:5]):  # 只预览前5页
                        try:
                            text = page.extract_text()
                            if text:
                                content += f"=== 第 {i + 1} 页 ===\n{text}\n\n"
                            else:
                                content += f"=== 第 {i + 1} 页 (无文本内容) ===\n\n"
                            self.progress_bar.setValue(i + 1)
                        except Exception as page_error:
                            content += f"=== 第 {i + 1} 页 (提取错误: {str(page_error)}) ===\n\n"

            if not content:
                content = "无法提取PDF文本内容"

            self.preview_text.setPlainText(content)
            self.status_label.setText(f"PDF预览: {os.path.basename(filepath)}")

        except Exception as e:
            self.preview_text.setPlainText(f"PDF加载失败: {str(e)}")
            self.status_label.setText("预览错误")

        finally:
            self.progress_bar.setVisible(False)

    def _show_pdf_image_preview(self, filepath):
        """图像预览模式"""
        self.status_label.setText("正在渲染PDF预览...")
        self.progress_bar.setVisible(True)
        self.progress_bar.setRange(0, 0)

        QTimer.singleShot(100, lambda: self._render_pdf_image_preview(filepath))

    def _render_pdf_image_preview(self, filepath):
        """渲染PDF图像预览"""
        try:
            # 创建临时目录
            temp_dir = tempfile.mkdtemp()

            # 渲染第一页为图像
            page_index = 0
            zoom = 1.0
            image_path = os.path.join(temp_dir, "preview.png")

            # 使用PyMuPDF渲染预览图
            with fitz.open(filepath) as doc:
                if len(doc) == 0:
                    raise Exception("空PDF文档")

                # 获取页面
                page = doc.load_page(page_index)

                # 设置渲染矩阵
                mat = fitz.Matrix(zoom, zoom)

                # 渲染页面为Pixmap
                pix = page.get_pixmap(matrix=mat, alpha=False)

                # 保存为PNG
                pix.save(image_path)

            # 显示图像预览
            pixmap = QPixmap(image_path)
            if not pixmap.isNull():
                self.preview_text.clear()
                self.preview_text.document().addResource(
                    QTextDocument.ImageResource,
                    QUrl("image://preview"),
                    pixmap
                )
                cursor = self.preview_text.textCursor()
                cursor.insertHtml(f'<center><img src="image://preview"></center>')
                self.status_label.setText(f"PDF图像预览: {os.path.basename(filepath)}")
            else:
                self.preview_text.setPlainText("无法渲染图像预览")
                self.status_label.setText("渲染失败")

        except Exception as e:
            self.preview_text.setPlainText(f"图像预览失败: {str(e)}\n已自动切换到文本预览模式")
            self.safe_mode_btn.setChecked(True)  # 自动切换到安全模式
            self._show_pdf_text_preview(filepath)  # 回退到文本预览

        finally:
            # 清理临时文件
            try:
                shutil.rmtree(temp_dir)
            except:
                pass

            self.progress_bar.setVisible(False)

    # Word文档预览
    def show_docx_preview(self, filepath):
        """预览Word文档"""
        self.status_label.setText("正在加载Word文档...")
        self.progress_bar.setVisible(True)
        self.progress_bar.setRange(0, 0)

        QTimer.singleShot(100, lambda: self._load_docx_content(filepath))

    def _load_docx_content(self, filepath):
        """加载Word内容"""
        try:
            doc = Document(filepath)
            content = ""
            paragraphs = doc.paragraphs[:100]  # 限制前100段
            self.progress_bar.setRange(0, len(paragraphs))

            for i, para in enumerate(paragraphs):
                content += para.text + "\n"
                if i % 10 == 0:  # 每10段更新一次进度
                    self.progress_bar.setValue(i)

            if not content:
                content = "无法提取Word文档内容"

            self.preview_text.setPlainText(content)
            self.status_label.setText(f"Word预览: {os.path.basename(filepath)}")

        except Exception as e:
            self.preview_text.setPlainText(f"Word文件加载失败: {str(e)}")
            self.status_label.setText("预览错误")

        finally:
            self.progress_bar.setVisible(False)

    # PPT文档预览
    def show_pptx_preview(self, filepath):
        """预览PPT文档"""
        self.status_label.setText("正在加载PPT文档...")
        self.progress_bar.setVisible(True)
        self.progress_bar.setRange(0, 0)

        QTimer.singleShot(100, lambda: self._load_pptx_content(filepath))

    def _load_pptx_content(self, filepath):
        """加载PPT内容"""
        try:
            prs = Presentation(filepath)
            content = ""
            slides = prs.slides[:10]  # 限制前10页
            self.progress_bar.setRange(0, len(slides))

            for i, slide in enumerate(slides):
                slide_title = f"幻灯片 {i + 1}"
                if slide.shapes.title:
                    slide_title = slide.shapes.title.text

                content += f"=== {slide_title} ===\n"

                # 提取内容
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape != slide.shapes.title:
                        slide_content.append(shape.text)

                content += "\n".join(slide_content) + "\n\n"
                self.progress_bar.setValue(i + 1)

            if not content:
                content = "无法提取PPT内容"

            self.preview_text.setPlainText(content)
            self.status_label.setText(f"PPT预览: {os.path.basename(filepath)}")

        except Exception as e:
            self.preview_text.setPlainText(f"PPT文件加载失败: {str(e)}")
            self.status_label.setText("预览错误")

        finally:
            self.progress_bar.setVisible(False)


class MainWindow(QMainWindow):
    """主应用窗口"""

    def __init__(self):
        super().__init__()
        self.setWindowTitle("大学生资料整理平台")
        self.setGeometry(100, 100, 1200, 800)

        # 数据结构
        self.documents = []
        self.knowledge_tree = []

        # 创建UI
        self.init_ui()

        # 预览窗口
        self.preview_window = PreviewWindow()

        # 添加上下文菜单
        self.setup_context_menu()

    def setup_context_menu(self):
        """添加上下文菜单"""
        # 文件列表的右键菜单
        self.file_list.setContextMenuPolicy(Qt.CustomContextMenu)
        self.file_list.customContextMenuRequested.connect(self.show_file_context_menu)

        # 知识树的右键菜单
        self.tree_widget.setContextMenuPolicy(Qt.CustomContextMenu)
        self.tree_widget.customContextMenuRequested.connect(self.show_tree_context_menu)

    def show_file_context_menu(self, position):
        """显示文件列表的右键菜单"""
        menu = QMenu()

        # 删除选中文件
        delete_action = QAction("删除选中文件", self)
        delete_action.triggered.connect(self.delete_selected_files)
        menu.addAction(delete_action)

        # 清空所有文件
        clear_action = QAction("清空所有文件", self)
        clear_action.triggered.connect(self.clear_all_files)
        menu.addAction(clear_action)

        menu.exec_(self.file_list.mapToGlobal(position))

    def show_tree_context_menu(self, position):
        """显示知识树的右键菜单"""
        menu = QMenu()

        # 删除选中节点
        delete_action = QAction("删除选中节点", self)
        delete_action.triggered.connect(self.delete_selected_node)
        menu.addAction(delete_action)

        # 编辑节点内容
        edit_action = QAction("编辑节点内容", self)
        edit_action.triggered.connect(self.edit_selected_node)
        menu.addAction(edit_action)

        menu.exec_(self.tree_widget.mapToGlobal(position))

    def delete_selected_files(self):
        """删除选中的文件"""
        selected_items = self.file_list.selectedItems()
        if not selected_items:
            QMessageBox.warning(self, "警告", "请先选择要删除的文件")
            return

        # 确认删除
        reply = QMessageBox.question(
            self,
            "确认删除",
            f"确定要删除选中的 {len(selected_items)} 个文件吗？",
            QMessageBox.Yes | QMessageBox.No
        )

        if reply == QMessageBox.Yes:
            # 从数据结构中移除
            for item in selected_items:
                file_name = item.text()
                # 找到对应的文档并从数据结构中移除
                self.documents = [doc for doc in self.documents if os.path.basename(doc['filepath']) != file_name]
                # 从列表中移除
                self.file_list.takeItem(self.file_list.row(item))

            self.update_status(f"已删除 {len(selected_items)} 个文件")
            # 清除知识树
            self.knowledge_tree = []
            self.tree_widget.clear()

    def clear_all_files(self):
        """清空所有文件"""
        if self.file_list.count() == 0:
            return

        # 确认清除
        reply = QMessageBox.question(
            self,
            "确认清除",
            "确定要清空所有文件吗？",
            QMessageBox.Yes | QMessageBox.No
        )

        if reply == QMessageBox.Yes:
            self.file_list.clear()
            self.documents = []
            self.knowledge_tree = []
            self.tree_widget.clear()
            self.update_status("已清空所有文件")

    def delete_selected_node(self):
        """删除选中的知识树节点"""
        selected_items = self.tree_widget.selectedItems()
        if not selected_items:
            QMessageBox.warning(self, "警告", "请先选择要删除的节点")
            return

        # 确认删除
        reply = QMessageBox.question(
            self,
            "确认删除",
            f"确定要删除选中的 {len(selected_items)} 个节点吗？",
            QMessageBox.Yes | QMessageBox.No
        )

        if reply == QMessageBox.Yes:
            # 从知识树中移除
            for item in selected_items:
                node_index = self.tree_widget.indexOfTopLevelItem(item)
                if 0 <= node_index < len(self.knowledge_tree):
                    del self.knowledge_tree[node_index]
                self.tree_widget.takeTopLevelItem(self.tree_widget.indexOfTopLevelItem(item))

            self.update_status(f"已删除 {len(selected_items)} 个节点")

    def edit_selected_node(self):
        """编辑选中的节点 - 简化版"""
        selected_items = self.tree_widget.selectedItems()
        if not selected_items or len(selected_items) > 1:
            QMessageBox.warning(self, "警告", "请选择单个节点进行编辑")
            return

        item = selected_items[0]
        node_data = item.data(0, Qt.UserRole)

        # 这里简化处理，实际应用中应该弹出编辑对话框
        QMessageBox.information(
            self,
            "节点编辑",
            f"实际应用中这里应该提供编辑界面\n当前节点标题: {node_data['title']}"
        )


    def init_ui(self):
        """初始化用户界面"""
        # 主布局
        main_widget = QWidget()
        main_layout = QHBoxLayout()

        # 左侧面板：文件操作
        left_panel = QVBoxLayout()

        # 文件列表
        self.file_list = QListWidget()
        self.file_list.setMaximumWidth(300)

        # 操作按钮
        btn_upload = QPushButton("上传文件")
        btn_upload.clicked.connect(self.upload_files)

        btn_generate = QPushButton("生成知识结构")
        btn_generate.clicked.connect(self.generate_knowledge_tree)

        # 导出选项
        export_layout = QHBoxLayout()
        self.export_combo = QComboBox()
        self.export_combo.addItems(["导出为Word", "导出为PDF"])
        btn_export = QPushButton("导出")
        btn_export.clicked.connect(self.export_document)
        export_layout.addWidget(self.export_combo)
        export_layout.addWidget(btn_export)

        left_panel.addWidget(QLabel("已上传文件:"))
        left_panel.addWidget(self.file_list)
        left_panel.addWidget(btn_upload)
        left_panel.addWidget(btn_generate)
        left_panel.addLayout(export_layout)

        # 右侧面板：知识树和预览
        right_panel = QVBoxLayout()

        # 知识结构树
        self.tree_widget = QTreeWidget()
        self.tree_widget.setHeaderLabel("知识结构")
        self.tree_widget.itemClicked.connect(self.show_preview)

        right_panel.addWidget(QLabel("知识结构:"))
        right_panel.addWidget(self.tree_widget)

        # 组合布局
        splitter = QSplitter(Qt.Horizontal)
        left_widget = QWidget()
        left_widget.setLayout(left_panel)
        right_widget = QWidget()
        right_widget.setLayout(right_panel)

        splitter.addWidget(left_widget)
        splitter.addWidget(right_widget)
        splitter.setSizes([300, 900])

        # 主窗口布局
        main_layout.addWidget(splitter)
        main_widget.setLayout(main_layout)
        self.setCentralWidget(main_widget)

    def upload_files(self):
        """文件上传功能"""
        files, _ = QFileDialog.getOpenFileNames(
            self,
            "选择学习资料",
            "",
            "文档文件 (*.pdf *.docx *.pptx)"
        )

        if not files:
            return

        for file_path in files:
            # 添加到文件列表
            self.file_list.addItem(os.path.basename(file_path))
            # 解析文档结构
            self.process_file(file_path)

    def process_file(self, file_path):
        """处理单个文件并存储数据结构"""
        processor = DocumentProcessor()
        ext = os.path.splitext(file_path)[1].lower()

        try:
            if ext == '.pdf':
                content = processor.extract_pdf_content(file_path)
            elif ext == '.docx':
                content = processor.extract_docx_content(file_path)
            elif ext == '.pptx':
                content = processor.extract_pptx_content(file_path)
            else:
                return

            content['filepath'] = file_path
            self.documents.append(content)
            self.update_status(f"已加载: {os.path.basename(file_path)}")
        except Exception as e:
            self.update_status(f"错误: {str(e)}")

    def generate_knowledge_tree(self):
        """生成知识结构树"""
        if not self.documents:
            self.update_status("请先上传文件")
            return

        processor = DocumentProcessor()
        self.knowledge_tree = processor.generate_knowledge_tree(self.documents)

        # 更新树状视图
        self.tree_widget.clear()
        for i, node in enumerate(self.knowledge_tree):
            item = QTreeWidgetItem([f"{i + 1}. {node['title']}"])
            item.setData(0, Qt.UserRole, node)
            self.tree_widget.addTopLevelItem(item)

        self.update_status("知识结构已生成")

    def show_preview(self, item, column):
        """显示文件预览"""
        node_data = item.data(0, Qt.UserRole)
        filepath = node_data.get('filepath', '')

        if not filepath:
            return

        ext = os.path.splitext(filepath)[1].lower()

        self.preview_window.show()
        try:
            if ext == '.pdf':
                self.preview_window.show_pdf_preview(filepath)
            elif ext == '.docx':
                self.preview_window.show_docx_preview(filepath)
            elif ext == '.pptx':
                self.preview_window.show_pptx_preview(filepath)
        except Exception as e:
            self.update_status(f"预览失败: {str(e)}")

    def export_document(self):
        """导出整理后的文档（优化过滤器默认选中状态）"""
        if not self.knowledge_tree:
            self.update_status("请先生成知识结构")
            return

        processor = DocumentProcessor()
        format_choice = self.export_combo.currentText()  # 获取下拉框选择的格式（"导出为Word"或"导出为PDF"）

        # 根据选择的格式，设置对应的过滤器和默认文件名
        if "Word" in format_choice:
            default_filename = "知识结构总结.docx"
            file_filter = "Word文档 (*.docx);;PDF文件 (*.pdf)"
            selected_filter_index = 0  # 默认选中第1个过滤器（Word）
        else:  # PDF格式
            default_filename = "知识结构总结.pdf"
            file_filter = "PDF文件 (*.pdf);;Word文档 (*.docx)"  # 交换顺序，让PDF在前面
            selected_filter_index = 0  # 默认选中第1个过滤器（PDF）

        # 获取文件对话框的选中过滤器（通过split分割过滤器字符串，取对应索引的过滤器）
        filters = file_filter.split(';;')
        selected_filter = filters[selected_filter_index]

        # 打开文件保存对话框，指定默认过滤器
        output_file, _ = QFileDialog.getSaveFileName(
            self,
            "保存整理文档",
            default_filename,
            file_filter,
            selected_filter  # 关键：指定默认选中的过滤器
        )

        if not output_file:
            return  # 用户取消保存

        try:
            if "Word" in format_choice:
                # 处理Word导出（确保扩展名正确）
                if not output_file.lower().endswith('.docx'):
                    output_file += '.docx'
                result = processor.export_to_word(self.knowledge_tree, output_file)
            else:
                # 处理PDF导出（先生成临时Word，再转换）
                temp_word_path = os.path.join(tempfile.gettempdir(), f"temp_pdf_conv_{os.getpid()}.docx")
                processor.export_to_word(self.knowledge_tree, temp_word_path)

                # 确保PDF扩展名正确
                if not output_file.lower().endswith('.pdf'):
                    output_file += '.pdf'
                result = processor.export_to_pdf(temp_word_path, output_file)

                # 清理临时文件
                if os.path.exists(temp_word_path):
                    os.remove(temp_word_path)

            if result and os.path.exists(result):
                self.update_status(f"导出成功: {os.path.basename(result)}")
                # 自动打开文件
                if sys.platform == "win32":
                    os.startfile(result)
                else:
                    opener = "open" if sys.platform == "darwin" else "xdg-open"
                    subprocess.call([opener, result])
            else:
                self.update_status("导出失败，未生成文件")
        except Exception as e:
            self.update_status(f"导出失败: {str(e)}")

    def update_status(self, message):
        """更新状态栏"""
        self.statusBar().showMessage(message, 5000)


# 确保所有类定义完整后再实例化
if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())